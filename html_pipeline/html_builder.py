from pathlib import Path
import io
import json
import re
from pptx import Presentation
from pptx.util import Inches

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)


COMPACT_STYLE = """
.slide {
  gap: 12px !important;
  padding: 16px 18px 14px 18px !important;
}
.header {
  gap: 6px !important;
  max-height: 104px !important;
}
.title {
  font-size: 30px !important;
  line-height: 1.15 !important;
}
.subtitle {
  font-size: 13px !important;
  line-height: 1.2 !important;
}
.card,
.summary-card,
.footer-card {
  padding: 14px 16px !important;
}
.card-inner {
  padding: 14px 16px !important;
}
.card-title,
.mini-title,
.side-title,
.panel-title {
  font-size: 19px !important;
  line-height: 1.18 !important;
  margin-bottom: 10px !important;
}
.step,
.stage,
.launch-box {
  padding: 12px !important;
}
.timeline-card {
  justify-content: flex-start !important;
  gap: 10px !important;
}
.flow,
.summary,
.decision-steps,
.interest-grid,
.top-right,
.bottom-right,
.timeline,
.main,
.platforms,
.tags,
.stage-tags {
  gap: 10px !important;
}
.metric-main {
  gap: 12px !important;
  margin: 8px 0 8px !important;
}
.num-block {
  padding: 12px 12px 10px !important;
}
.tag,
.pill,
.platform {
  font-size: 11px !important;
  padding: 5px 9px !important;
}
.footer,
.summary-card,
.footer-card {
  height: 84px !important;
  min-height: 84px !important;
}
.summary-quote,
.footer-key {
  font-size: 22px !important;
  line-height: 1.15 !important;
}
.summary-sub,
.footer-text,
.metric-desc,
.step-text,
.stage-desc,
.panel-desc,
.launch-text,
.accept-line,
.lead,
.box-sub,
.metric-sub,
.mini-text {
  font-size: 13px !important;
  line-height: 1.35 !important;
}
.metric,
.metric-big,
.metric.small,
.launch-date,
.num {
  font-size: 32px !important;
  line-height: 1 !important;
}
.step-num,
.icon {
  transform: scale(0.92);
  transform-origin: center;
}
"""


TOC_SAFE_STYLE = """
.slide {
  gap: 14px !important;
}
.left-card {
  justify-content: flex-start !important;
  gap: 16px !important;
}
.flow-card {
  gap: 12px !important;
}
.steps {
  flex: 1 !important;
  gap: 12px !important;
}
.flow-fill {
  display: grid !important;
  grid-template-columns: repeat(3, 1fr) !important;
  gap: 12px !important;
  margin-top: 2px !important;
}
.fill-card {
  border-radius: 12px !important;
  padding: 12px 12px 10px !important;
  background: rgba(255,255,255,0.03) !important;
  border: 1px solid rgba(122, 137, 165, 0.14) !important;
}
.fill-label {
  font-size: 12px !important;
  font-weight: 800 !important;
  color: #8FA3BF !important;
  margin-bottom: 6px !important;
  letter-spacing: 0.3px !important;
  text-transform: uppercase !important;
}
.fill-text {
  font-size: 13px !important;
  line-height: 1.4 !important;
  color: #C5CFDC !important;
}
"""


STEP_CARD_SAFE_STYLE = """
.steps,
.process-grid,
.process-wrap {
  gap: 8px !important;
}
.step {
  grid-template-columns: 40px 1fr !important;
  gap: 8px !important;
  padding: 10px !important;
  min-width: 0 !important;
}
.step-num,
.step-no {
  width: 40px !important;
  height: 40px !important;
  font-size: 18px !important;
  border-radius: 10px !important;
  flex: 0 0 40px !important;
}
.step-body,
.step-head {
  gap: 3px !important;
  min-width: 0 !important;
}
.step-title,
.step h4 {
  font-size: 16px !important;
  line-height: 1.12 !important;
  overflow-wrap: anywhere !important;
  word-break: break-word !important;
}
.step-text,
.step-desc,
.step p,
.card-sub,
.mini-note {
  font-size: 12px !important;
  line-height: 1.26 !important;
  overflow-wrap: anywhere !important;
  word-break: break-word !important;
}
.step-note,
.step-time {
  font-size: 11px !important;
  line-height: 1.15 !important;
}
.split,
.mini-tags,
.step-tags {
  gap: 6px !important;
}
.mini-tag,
.pill-green,
.pill-red,
.step-tag {
  font-size: 11px !important;
  padding: 4px 8px !important;
}
"""


CONCLUSION_SAFE_STYLE = """
.slide {
  grid-template-rows: 92px 1fr 84px !important;
  gap: 12px !important;
  padding: 16px 18px 16px 18px !important;
}
.main {
  grid-template-columns: 1.08fr 0.92fr !important;
  gap: 16px !important;
}
.card {
  padding: 18px !important;
}
.mid-bridge {
  padding: 12px 14px !important;
  gap: 10px !important;
}
.bridge-label {
  font-size: 12px !important;
  padding: 5px 9px !important;
}
.bridge-item {
  font-size: 13px !important;
  line-height: 1.3 !important;
}
.arrow {
  font-size: 16px !important;
}
.steps {
  gap: 10px !important;
}
.step {
  grid-template-columns: 44px 1fr !important;
  gap: 10px !important;
  padding: 12px !important;
}
.step-num {
  font-size: 22px !important;
  border-radius: 10px !important;
}
.step-body {
  gap: 4px !important;
}
.step-title {
  font-size: 17px !important;
  line-height: 1.15 !important;
}
.step-desc {
  font-size: 13px !important;
  line-height: 1.32 !important;
}
.step-note {
  font-size: 11px !important;
}
.split-result {
  padding: 12px 14px !important;
  gap: 10px !important;
}
.result-box {
  padding: 10px 12px !important;
  gap: 5px !important;
}
.result-box strong {
  font-size: 15px !important;
  line-height: 1.15 !important;
}
.result-box p {
  font-size: 12px !important;
  line-height: 1.28 !important;
}
.footer {
  height: 84px !important;
}
.summary-card {
  padding: 12px 16px !important;
  gap: 14px !important;
}
.summary-pill {
  font-size: 11px !important;
  padding: 7px 10px !important;
}
.summary-text {
  font-size: 16px !important;
  line-height: 1.25 !important;
}
.summary-actions {
  gap: 8px !important;
}
.choice {
  font-size: 11px !important;
  padding: 8px 10px !important;
}
"""


SUMMARY_SAFE_STYLE = """
.slide {
  grid-template-rows: auto 1fr 84px !important;
  gap: 12px !important;
  padding: 16px 18px 14px 18px !important;
}
.header,
.title-wrap {
  gap: 6px !important;
}
.title {
  font-size: 30px !important;
  line-height: 1.15 !important;
}
.subtitle {
  font-size: 13px !important;
  line-height: 1.2 !important;
}
.main {
  grid-template-columns: minmax(0, 1.12fr) minmax(0, 0.88fr) !important;
  gap: 12px !important;
}
.mid-card {
  display: none !important;
}
.card,
.card-inner,
.summary-card,
.footer-card {
  padding: 14px 16px !important;
}
.card-title,
.panel-title,
.side-title,
.mini-title {
  font-size: 19px !important;
  line-height: 1.18 !important;
  margin-bottom: 10px !important;
}
.timeline-card {
  justify-content: flex-start !important;
  gap: 10px !important;
}
.interest-grid,
.tags,
.platforms,
.summary,
.decision-steps,
.timeline {
  gap: 10px !important;
}
.metric-main {
  gap: 12px !important;
  margin: 8px 0 8px !important;
}
.num-block {
  padding: 12px 12px 10px !important;
}
.step,
.stage,
.launch-box,
.accept-line {
  padding: 12px !important;
}
.step:nth-child(n+3) {
  display: none !important;
}
.metric-card .accept-line {
  display: none !important;
}
.footer,
.summary-card,
.footer-card {
  height: 84px !important;
  min-height: 84px !important;
}
.summary-quote,
.footer-key {
  font-size: 22px !important;
  line-height: 1.15 !important;
}
.tag,
.pill,
.platform {
  font-size: 11px !important;
  padding: 5px 9px !important;
}
.metric,
.metric-big,
.metric.small,
.launch-date,
.num {
  font-size: 32px !important;
  line-height: 1 !important;
}
.panel-desc,
.step p,
.metric-desc,
.metric-sub,
.launch-text,
.summary-sub,
.footer-text,
.accept-line,
.lead,
.box-sub,
.mini-text {
  font-size: 13px !important;
  line-height: 1.35 !important;
}
"""


TIMELINE_SAFE_STYLE = """
.slide {
  grid-template-rows: 88px 1fr 84px !important;
  gap: 12px !important;
  padding: 16px 18px 14px 18px !important;
}
.header,
.title-wrap {
  gap: 6px !important;
}
h1,
.title {
  font-size: 29px !important;
  line-height: 1.15 !important;
  max-width: 900px !important;
}
.subtitle {
  font-size: 13px !important;
  line-height: 1.2 !important;
}
.header-badge {
  min-width: 168px !important;
  padding: 10px 12px !important;
}
.main {
  grid-template-columns: 1.04fr 0.96fr !important;
  gap: 12px !important;
}
.card,
.metric-card,
.node-card,
.footer-card {
  padding: 16px !important;
}
.card-title,
.metric-title,
.node-name,
.step-title {
  font-size: 18px !important;
  line-height: 1.18 !important;
}
.timeline {
  grid-template-columns: repeat(4, 1fr) !important;
  gap: 8px !important;
}
.timeline .step:nth-child(n+5) {
  display: none !important;
}
.step {
  padding-top: 12px !important;
  min-width: 0 !important;
}
.step-head {
  gap: 6px !important;
  margin-bottom: 8px !important;
  min-width: 0 !important;
}
.step-head .num,
.num {
  width: 22px !important;
  height: 22px !important;
  font-size: 11px !important;
  flex: 0 0 22px !important;
}
.step-title {
  font-size: 14px !important;
  line-height: 1.15 !important;
  overflow-wrap: anywhere !important;
  word-break: break-word !important;
}
.step-body {
  padding-left: 28px !important;
  gap: 4px !important;
  min-width: 0 !important;
}
.step-time {
  font-size: 11px !important;
  line-height: 1.15 !important;
  overflow-wrap: anywhere !important;
  word-break: break-word !important;
}
.step-text,
.node-desc,
.mini-summary p,
.footer-right,
.metric .note {
  font-size: 12px !important;
  line-height: 1.3 !important;
  overflow-wrap: anywhere !important;
  word-break: break-word !important;
}
.right-col {
  grid-template-rows: 1fr 152px !important;
  gap: 12px !important;
}
.highlight-block,
.node-grid,
.tags {
  gap: 10px !important;
}
.metric {
  padding: 14px !important;
}
.metric .number {
  font-size: 32px !important;
  margin-bottom: 4px !important;
}
.footer-card {
  grid-template-columns: 170px 1fr !important;
  gap: 14px !important;
  padding: 12px 16px !important;
}
.footer-left .big {
  font-size: 24px !important;
}
"""


DENSE_CARD_SAFE_STYLE = """
.slide {
  gap: 12px !important;
  padding: 16px 18px 14px 18px !important;
  grid-template-rows: auto 1fr 88px !important;
}
.header,
.title-wrap {
  gap: 6px !important;
}
h1,
.title {
  font-size: 29px !important;
  line-height: 1.14 !important;
}
.subtitle {
  font-size: 13px !important;
  line-height: 1.2 !important;
}
.main,
.side,
.right-col,
.right-grid,
.right-panel,
.stat-grid,
.highlight-block,
.logic-row,
.timeline,
.timeline-wrap,
.node-grid,
.tags,
.mini-tags,
.tag-row,
.summary-tags,
.footer-tags,
.bullet-list,
.support-list,
.process-grid,
.metrics-grid,
.metrics,
.scene-top,
.bottom {
  gap: 8px !important;
}
.card,
.metric-card,
.stats-card,
.evidence-card,
.talk-card,
.node-card,
.footer-card,
.panel,
.memory-card,
.stage,
.summary-card,
.metric-box,
.support-item,
.time-box,
.scene-item,
.policy-box {
  padding: 12px !important;
}
.card-title,
.stats-title,
.metric-title,
.hero-title,
.big-conclusion,
.node-name,
.logic-name,
.step-title,
.memory-title,
.stage-name,
.footer-title,
.section-title,
.support-title {
  font-size: 16px !important;
  line-height: 1.12 !important;
}
.card-label,
.eyebrow,
.source-tag,
.badge-text,
.metric-badge,
.tag,
.mini-tag,
.platform,
.pill,
.memory-kicker,
.stage-label,
.footer-tag,
.step-tag {
  font-size: 10px !important;
}
.hero-desc,
.short-text,
.logic-text,
.node-text,
.node-desc,
.metric-desc,
.metric-sub,
.step-text,
.step-desc,
.note,
.desc,
.support,
.footer-right,
.summary-sub,
.footer-text,
.lead,
.stage-desc,
.summary-text,
.metric-text,
.metric-caption,
.bullet,
.support-text,
.keyline-text,
.card-sub,
.timeline-summary,
.metric-note,
.scene-text,
.policy-title,
.policy-note,
.card-note,
.mini-note {
  font-size: 11px !important;
  line-height: 1.24 !important;
}
.stat,
.metric,
.num-block,
.logic-item,
.node,
.step,
.fromto,
.stage,
.support-item,
.time-box,
.metric-box,
.scene-item,
.policy-box {
  padding: 10px !important;
}
.stat-value,
.stat-num,
.metric .number,
.metric-big,
.metric,
.big,
.stage-year,
.year,
.metric-main,
.big-number,
.metric-number,
.policy-year .big {
  font-size: 26px !important;
  line-height: 1 !important;
}
.step,
.node,
.logic-item,
.card,
.metric-card,
.step-title,
.step-text,
.node-name,
.node-desc,
.logic-name,
.logic-text,
.memory-title,
.footer-right,
.stage,
.stage-name,
.stage-desc,
.summary-text,
.metric-text,
.metric-caption,
.support-text,
.footer-text,
.lead,
.keyline-text,
.card-sub,
.timeline-summary,
.metric-note,
.scene-text,
.policy-title,
.mini-note {
  overflow-wrap: anywhere !important;
  word-break: break-word !important;
}
.step-head .num,
.num,
.step-no {
  width: 22px !important;
  height: 22px !important;
  font-size: 11px !important;
  flex: 0 0 22px !important;
}
.step-head {
  gap: 8px !important;
}
.step-body {
  padding-left: 28px !important;
  gap: 4px !important;
}
.timeline {
  grid-template-columns: repeat(4, 1fr) !important;
}
.timeline .step:nth-child(n+5),
.timeline-wrap .node:nth-child(n+5),
.timeline .node:nth-child(n+5),
.logic-row .logic-item:nth-child(n+4),
.node-grid .node:nth-child(n+5),
.tag-row .tag:nth-child(n+4),
.summary-tags .tag:nth-child(n+4),
.footer-tags .footer-tag:nth-child(n+3),
.tags .tag:nth-child(n+4),
.bullet-list .bullet:nth-child(n+3),
.support-list .support-item:nth-child(n+3),
.metrics-grid .metric-box:nth-child(n+3),
.metrics .metric:nth-child(n+3),
.scene-top .scene-item:nth-child(n+4) {
  display: none !important;
}
.footer,
.summary-card,
.footer-card {
  height: 88px !important;
  min-height: 88px !important;
}
"""


def _inspect_html_layout(page) -> dict:
    return page.evaluate(
        """
        () => {
          const slide = document.querySelector('.slide');
          const header = document.querySelector('.header');
          const main = document.querySelector('.main');
          const footer = document.querySelector('.footer');

          const overflowSelectors = [
            '.card', '.stage', '.step', '.summary-card', '.footer-card',
            '.launch-box', '.summary', '.flow', '.node', '.metric', '.metric-box',
            '.scene-item', '.policy-box', '.timeline-summary', '.card-sub', '.metric-note'
          ];

          const getRect = (el) => {
            if (!el) return null;
            const r = el.getBoundingClientRect();
            return {
              top: Math.round(r.top),
              bottom: Math.round(r.bottom),
              left: Math.round(r.left),
              right: Math.round(r.right),
              width: Math.round(r.width),
              height: Math.round(r.height),
            };
          };

          const collectOverflow = () => {
            const items = [];
            for (const selector of overflowSelectors) {
              document.querySelectorAll(selector).forEach((el, index) => {
                const sh = el.scrollHeight;
                const ch = el.clientHeight;
                const sw = el.scrollWidth;
                const cw = el.clientWidth;
                if (sh > ch + 2 || sw > cw + 2) {
                  items.push({
                    selector,
                    index,
                    scrollHeight: sh,
                    clientHeight: ch,
                    scrollWidth: sw,
                    clientWidth: cw,
                    text: (el.innerText || '').replace(/\\s+/g, ' ').slice(0, 80),
                  });
                }
              });
            }
            return items;
          };

          const slideRect = getRect(slide);
          const headerRect = getRect(header);
          const mainRect = getRect(main);
          const footerRect = getRect(footer);
          const overlap = Boolean(mainRect && footerRect && mainRect.bottom > footerRect.top);
          const slideOverflow = slide ? slide.scrollHeight > slide.clientHeight + 2 : false;

          const boundarySelectors = [
            '.title', '.subtitle', '.summary-card', '.footer-card',
            '.card-title', '.card-sub', '.footer-text', '.timeline-summary', '.policy-title'
          ];
          const boundaryIssues = [];
          if (slideRect) {
            for (const selector of boundarySelectors) {
              document.querySelectorAll(selector).forEach((el, index) => {
                const rect = getRect(el);
                if (!rect) return;
                if (rect.right > slideRect.right + 1 || rect.bottom > slideRect.bottom + 1) {
                  boundaryIssues.push({ selector, index, rect });
                }
              });
            }
          }

          const summaryHeavy = Boolean(
            footer
            && main
            && document.querySelector('.decision-steps')
            && document.querySelector('.metric-card')
            && window.getComputedStyle(main).gridTemplateColumns.split(' ').length >= 3
          );

          const timelineHeavy = Boolean(
            (document.querySelector('.timeline-wrap') && document.querySelectorAll('.timeline-wrap .node').length >= 4)
            || (
              document.querySelector('.timeline')
              && (
                document.querySelectorAll('.timeline .step').length >= 4
                || document.querySelectorAll('.timeline .node').length >= 4
              )
              && document.querySelector('.right-col, .right-grid, .metrics, .metrics-grid, .bottom')
            )
          );

          const denseCardHeavy = Boolean(
            document.querySelectorAll('.card, .metric-card, .stats-card, .evidence-card, .talk-card, .scene-card').length >= 2
            && (
              document.querySelector('.logic-row')
              || document.querySelector('.stat-grid')
              || document.querySelector('.node-grid')
              || document.querySelector('.metrics-grid')
              || document.querySelector('.metrics')
              || document.querySelector('.scene-top')
              || document.querySelector('.policy-box')
              || document.querySelector('.bottom')
            )
          );

          const tocSparse = Boolean(
            document.querySelector('.steps')
            && document.querySelector('.flow-card')
            && document.querySelector('.left-card')
            && !document.querySelector('.flow-fill')
          );

          const conclusionHeavy = Boolean(
            document.querySelector('.decision-card')
            && document.querySelector('.split-result')
            && document.querySelector('.summary-card')
          );

          const stepCardHeavy = Boolean(
            document.querySelector('.steps')
            && document.querySelectorAll('.steps .step').length >= 3
          );

          return {
            slide: slide ? {
              scrollHeight: slide.scrollHeight,
              clientHeight: slide.clientHeight,
              scrollWidth: slide.scrollWidth,
              clientWidth: slide.clientWidth,
              rect: slideRect,
            } : null,
            header: header ? { rect: headerRect } : null,
            main: main ? { rect: mainRect } : null,
            footer: footer ? { rect: footerRect } : null,
            overlap,
            slideOverflow,
            overflowItems: collectOverflow(),
            boundaryIssues,
            summaryHeavy,
            timelineHeavy,
            denseCardHeavy,
            tocSparse,
            conclusionHeavy,
            stepCardHeavy,
          };
        }
        """
    )


def _summarize_layout_issues(report: dict) -> list[str]:
    issues = []
    slide = report.get("slide") or {}
    header = report.get("header") or {}

    if report.get("slideOverflow"):
        issues.append(
            f"slide 高度超限：scrollHeight={slide.get('scrollHeight')} > clientHeight={slide.get('clientHeight')}"
        )

    if report.get("overlap"):
        main_rect = (report.get("main") or {}).get("rect") or {}
        footer_rect = (report.get("footer") or {}).get("rect") or {}
        issues.append(
            f"main/footer 重叠：main.bottom={main_rect.get('bottom')} > footer.top={footer_rect.get('top')}"
        )

    header_rect = header.get("rect") or {}
    if header_rect.get("height", 0) > 110:
        issues.append(f"header 过高：height={header_rect.get('height')}px")

    for item in (report.get("overflowItems") or [])[:5]:
        if item.get("scrollHeight", 0) > item.get("clientHeight", 0) + 2:
            issues.append(
                f"{item['selector']}#{item['index']} 垂直溢出：{item['scrollHeight']}>{item['clientHeight']} 文本={item['text']}"
            )
        elif item.get("scrollWidth", 0) > item.get("clientWidth", 0) + 2:
            issues.append(
                f"{item['selector']}#{item['index']} 水平溢出：{item['scrollWidth']}>{item['clientWidth']} 文本={item['text']}"
            )

    for item in (report.get("boundaryIssues") or [])[:3]:
        rect = item.get("rect") or {}
        issues.append(
            f"{item['selector']}#{item['index']} 超出 slide 边界：right={rect.get('right')}, bottom={rect.get('bottom')}"
        )

    return issues


def _should_regenerate(report: dict) -> bool:
    overflow_items = report.get("overflowItems") or []
    severe_overflow_count = sum(
        1
        for item in overflow_items
        if (item.get("scrollHeight", 0) - item.get("clientHeight", 0) > 14)
        or (item.get("scrollWidth", 0) - item.get("clientWidth", 0) > 14)
    )
    total_overflow_count = len(overflow_items)
    header_height = ((report.get("header") or {}).get("rect") or {}).get("height", 0)
    slide_delta = (report.get("slide") or {}).get("scrollHeight", 0) - (report.get("slide") or {}).get("clientHeight", 0)

    high_risk_dense_page = report.get("timelineHeavy") or report.get("denseCardHeavy")

    return bool(
        report.get("overlap")
        or severe_overflow_count >= 1
        or total_overflow_count >= 3
        or header_height > 118
        or slide_delta > 12
        or len(report.get("boundaryIssues") or []) >= 2
        or (high_risk_dense_page and total_overflow_count >= 1)
    )


def _apply_compact_mode(page) -> None:
    page.add_style_tag(content=COMPACT_STYLE)


def _apply_summary_safe_mode(page) -> None:
    page.add_style_tag(content=SUMMARY_SAFE_STYLE)


def _apply_timeline_safe_mode(page) -> None:
    page.add_style_tag(content=TIMELINE_SAFE_STYLE)


def _apply_dense_card_safe_mode(page) -> None:
    page.add_style_tag(content=DENSE_CARD_SAFE_STYLE)


def _apply_toc_safe_mode(page) -> None:
    page.add_style_tag(content=TOC_SAFE_STYLE)


def _apply_conclusion_safe_mode(page) -> None:
    page.add_style_tag(content=CONCLUSION_SAFE_STYLE)


def _apply_step_card_safe_mode(page) -> None:
    page.add_style_tag(content=STEP_CARD_SAFE_STYLE)


def _persist_style(html_path: Path, original_html: str, style_id: str, style_content: str) -> str:
    style_tag = f"<style id=\"{style_id}\">\n{style_content}\n</style>"
    if f'id="{style_id}"' in original_html:
        updated_html = re.sub(
            rf'<style id="{re.escape(style_id)}">.*?</style>',
            style_tag,
            original_html,
            count=1,
            flags=re.DOTALL,
        )
    elif "</head>" in original_html:
        updated_html = original_html.replace("</head>", f"{style_tag}\n</head>", 1)
    else:
        updated_html = original_html + "\n" + style_tag
    html_path.write_text(updated_html, encoding="utf-8")
    return updated_html


def _persist_compact_html(html_path: Path, original_html: str) -> str:
    return _persist_style(html_path, original_html, "claude-compact-style", COMPACT_STYLE)


def _persist_summary_safe_html(html_path: Path, original_html: str) -> str:
    return _persist_style(html_path, original_html, "claude-summary-safe-style", SUMMARY_SAFE_STYLE)


def _persist_timeline_safe_html(html_path: Path, original_html: str) -> str:
    return _persist_style(html_path, original_html, "claude-timeline-safe-style", TIMELINE_SAFE_STYLE)


def _persist_dense_card_safe_html(html_path: Path, original_html: str) -> str:
    return _persist_style(html_path, original_html, "claude-dense-card-safe-style", DENSE_CARD_SAFE_STYLE)


def _persist_toc_safe_html(html_path: Path, original_html: str) -> str:
    return _persist_style(html_path, original_html, "claude-toc-safe-style", TOC_SAFE_STYLE)


def _persist_conclusion_safe_html(html_path: Path, original_html: str) -> str:
    return _persist_style(html_path, original_html, "claude-conclusion-safe-style", CONCLUSION_SAFE_STYLE)


def _persist_step_card_safe_html(html_path: Path, original_html: str) -> str:
    return _persist_style(html_path, original_html, "claude-step-card-safe-style", STEP_CARD_SAFE_STYLE)


def render_html_with_validation(html_path: Path) -> tuple[bytes, dict]:
    from playwright.sync_api import sync_playwright

    html_content = html_path.read_text(encoding="utf-8")
    persisted_compact = False
    persisted_summary_safe = False
    persisted_timeline_safe = False
    persisted_dense_card_safe = False
    persisted_toc_safe = False
    persisted_conclusion_safe = False
    persisted_step_card_safe = False
    summary_safe_applied = False
    timeline_safe_applied = False
    dense_card_safe_applied = False
    toc_safe_applied = False
    conclusion_safe_applied = False
    step_card_safe_applied = False
    with sync_playwright() as p:
        browser = p.chromium.launch()
        page = browser.new_page(viewport={"width": 1280, "height": 720})
        page.set_content(html_content, wait_until="networkidle")

        initial_report = _inspect_html_layout(page)
        initial_issues = _summarize_layout_issues(initial_report)

        final_report = initial_report
        final_issues = initial_issues

        if initial_report.get("tocSparse"):
            _apply_toc_safe_mode(page)
            toc_safe_applied = True
            persisted_toc_safe = True
            final_report = _inspect_html_layout(page)
            final_issues = _summarize_layout_issues(final_report)
            html_content = _persist_toc_safe_html(html_path, html_content)

        if final_report.get("timelineHeavy") and (final_report.get("overlap") or len(final_issues) >= 3):
            _apply_timeline_safe_mode(page)
            timeline_safe_applied = True
            persisted_timeline_safe = True
            final_report = _inspect_html_layout(page)
            final_issues = _summarize_layout_issues(final_report)
            html_content = _persist_timeline_safe_html(html_path, html_content)

        if final_report.get("denseCardHeavy") and len(final_issues) >= 1:
            _apply_dense_card_safe_mode(page)
            dense_card_safe_applied = True
            persisted_dense_card_safe = True
            final_report = _inspect_html_layout(page)
            final_issues = _summarize_layout_issues(final_report)
            html_content = _persist_dense_card_safe_html(html_path, html_content)

        if final_report.get("summaryHeavy") and (final_report.get("overlap") or len(final_issues) >= 2):
            _apply_summary_safe_mode(page)
            summary_safe_applied = True
            persisted_summary_safe = True
            final_report = _inspect_html_layout(page)
            final_issues = _summarize_layout_issues(final_report)
            html_content = _persist_summary_safe_html(html_path, html_content)

        if final_report.get("conclusionHeavy") and len(final_issues) >= 1:
            _apply_conclusion_safe_mode(page)
            conclusion_safe_applied = True
            persisted_conclusion_safe = True
            final_report = _inspect_html_layout(page)
            final_issues = _summarize_layout_issues(final_report)
            html_content = _persist_conclusion_safe_html(html_path, html_content)

        if final_report.get("stepCardHeavy") and len(final_issues) >= 1:
            _apply_step_card_safe_mode(page)
            step_card_safe_applied = True
            persisted_step_card_safe = True
            final_report = _inspect_html_layout(page)
            final_issues = _summarize_layout_issues(final_report)
            html_content = _persist_step_card_safe_html(html_path, html_content)

        compact_applied = False
        if final_issues and not _should_regenerate(final_report):
            _apply_compact_mode(page)
            compact_applied = True
            persisted_compact = True
            final_report = _inspect_html_layout(page)
            final_issues = _summarize_layout_issues(final_report)
            html_content = _persist_compact_html(html_path, html_content)

        png = page.screenshot(type="png", clip={
            "x": 0, "y": 0, "width": 1280, "height": 720,
        })
        browser.close()

    status = "pass"
    if final_issues:
        status = "regenerate" if _should_regenerate(final_report) else "compact_pass"

    return png, {
        "status": status,
        "compact_applied": compact_applied,
        "summary_safe_applied": summary_safe_applied,
        "timeline_safe_applied": timeline_safe_applied,
        "dense_card_safe_applied": dense_card_safe_applied,
        "toc_safe_applied": toc_safe_applied,
        "conclusion_safe_applied": conclusion_safe_applied,
        "step_card_safe_applied": step_card_safe_applied,
        "persisted_compact": persisted_compact,
        "persisted_summary_safe": persisted_summary_safe,
        "persisted_timeline_safe": persisted_timeline_safe,
        "persisted_dense_card_safe": persisted_dense_card_safe,
        "persisted_toc_safe": persisted_toc_safe,
        "persisted_conclusion_safe": persisted_conclusion_safe,
        "persisted_step_card_safe": persisted_step_card_safe,
        "initial_issues": initial_issues,
        "final_issues": final_issues,
        "initial_report": initial_report,
        "final_report": final_report,
    }


def render_html_screenshot(html_path: Path) -> bytes:
    png, _ = render_html_with_validation(html_path)
    return png


def save_html_screenshot(html_path: Path, image_path: Path) -> Path:
    png_bytes = render_html_screenshot(html_path)
    image_path.write_bytes(png_bytes)
    return image_path


def write_slide_status(out_dir: Path, slide_status: dict) -> Path:
    status_path = out_dir / "slide-status.json"
    status_path.write_text(
        json.dumps({"slides": slide_status}, ensure_ascii=False, indent=2),
        encoding="utf-8",
    )
    return status_path


def html_to_png_bytes(html_path: Path) -> bytes:
    """用 playwright 将 HTML 文件渲染为 1280×720 PNG。"""
    png, _ = render_html_with_validation(html_path)
    return png


def build_pptx(html_dir: Path, output_path: Path) -> Path:
    """将 HTML 目录中的所有 HTML 文件转换为 PPTX。"""
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    blank_layout = prs.slide_layouts[6]

    html_files = sorted(html_dir.glob("*.html"))
    if not html_files:
        raise ValueError(f"HTML 目录为空: {html_dir}")

    for html_path in html_files:
        print(f"  插入: {html_path.name}")
        slide = prs.slides.add_slide(blank_layout)
        try:
            png_bytes, report = render_html_with_validation(html_path)
            if report.get("toc_safe_applied"):
                print("    [检查] 检测到目录页留白过大，已应用 toc-safe mode")
            if report.get("timeline_safe_applied"):
                print("    [检查] 检测到时间线高风险布局，已应用 timeline-safe mode")
            if report.get("dense_card_safe_applied"):
                print("    [检查] 检测到高密度卡片布局，已应用 dense-card-safe mode")
            if report.get("summary_safe_applied"):
                print("    [检查] 检测到总结型高风险布局，已应用 summary-safe mode")
            if report.get("conclusion_safe_applied"):
                print("    [检查] 检测到结论页过挤布局，已应用 conclusion-safe mode")
            if report.get("step_card_safe_applied"):
                print("    [检查] 检测到步骤卡纵向超限，已应用 step-card-safe mode")
            if report["compact_applied"]:
                print("    [检查] 检测到轻微超限，已应用紧凑模式")
            if report.get("persisted_toc_safe"):
                print("    [检查] 已将 toc-safe 样式回写到 HTML 文件")
            if report.get("persisted_timeline_safe"):
                print("    [检查] 已将 timeline-safe 样式回写到 HTML 文件")
            if report.get("persisted_dense_card_safe"):
                print("    [检查] 已将 dense-card-safe 样式回写到 HTML 文件")
            if report.get("persisted_summary_safe"):
                print("    [检查] 已将 summary-safe 样式回写到 HTML 文件")
            if report.get("persisted_conclusion_safe"):
                print("    [检查] 已将 conclusion-safe 样式回写到 HTML 文件")
            if report.get("persisted_step_card_safe"):
                print("    [检查] 已将 step-card-safe 样式回写到 HTML 文件")
            if report.get("persisted_compact"):
                print("    [检查] 已将紧凑版样式回写到 HTML 文件")
            if report["final_issues"]:
                print(f"    [检查] 仍有 {len(report['final_issues'])} 个布局问题")
            slide.shapes.add_picture(
                io.BytesIO(png_bytes), left=0, top=0,
                width=SLIDE_W, height=SLIDE_H,
            )
        except Exception as e:
            print(f"  [警告] {html_path.name} 失败: {e}，跳过")

    prs.save(str(output_path))
    print(f"PPT 已保存: {output_path}")
    return output_path
