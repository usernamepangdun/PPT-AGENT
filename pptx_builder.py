from pathlib import Path
import io
from pptx import Presentation
from pptx.util import Inches

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)


def svg_to_png_bytes(svg_path: Path) -> bytes:
    from playwright.sync_api import sync_playwright
    svg_content = svg_path.read_text(encoding="utf-8")
    html = f"""<!DOCTYPE html>
<html><head><meta charset="utf-8">
<style>*{{margin:0;padding:0;background:#fff}}</style></head>
<body>{svg_content}</body></html>"""
    with sync_playwright() as p:
        browser = p.chromium.launch()
        page = browser.new_page(viewport={"width": 1280, "height": 720})
        page.set_content(html, wait_until="networkidle")
        png = page.locator("svg").first.screenshot(type="png")
        browser.close()
    return png


def build_pptx(svg_dir: Path, output_path: Path) -> Path:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    blank_layout = prs.slide_layouts[6]

    svg_files = sorted(svg_dir.glob("*.svg"))
    if not svg_files:
        raise ValueError(f"SVG 目录为空: {svg_dir}")

    for svg_path in svg_files:
        print(f"  插入: {svg_path.name}")
        slide = prs.slides.add_slide(blank_layout)
        try:
            png_bytes = svg_to_png_bytes(svg_path)
            slide.shapes.add_picture(
                io.BytesIO(png_bytes), left=0, top=0,
                width=SLIDE_W, height=SLIDE_H
            )
        except Exception as e:
            print(f"  [警告] {svg_path.name} 失败: {e}，跳过")

    prs.save(str(output_path))
    print(f"PPT 已保存: {output_path}")
    return output_path
